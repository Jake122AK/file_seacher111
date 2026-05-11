r"""
FileSearcher — Windows用 全文検索ツール
========================================
- 指定フォルダ（ローカル / ネットワーク \\server\share）配下を再帰スキャン
- テキスト & Office (.docx/.xlsx/.pptx) & PDF を抽出
- SQLite FTS5 でインデックス化 → 2回目以降の検索は一瞬
- ファイル名 / 中身の両方にヒット、該当箇所のスニペットを表示
- 増分インデックス（mtime比較で変更分だけ再処理）

実行: python file_searcher.py
"""

from __future__ import annotations

import os
import sys
import re
import hashlib
import sqlite3
import subprocess
import threading
import queue
import time
import json
import traceback
from concurrent.futures import ThreadPoolExecutor, FIRST_COMPLETED, wait
from pathlib import Path
from datetime import datetime

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext

# ---------------------------------------------------------------------------
# 設定
# ---------------------------------------------------------------------------

APP_NAME = "FileSearcher"
APP_DIR = Path.home() / ".file_searcher"
APP_DIR.mkdir(exist_ok=True)
SHARD_DIR = APP_DIR / "shards"                # シャードDB群の格納場所
LEGACY_DB = APP_DIR / "index.db"              # 旧単一DBファイル（マイグレーション元）
CONFIG_PATH = APP_DIR / "config.json"
NUM_SHARDS = 16                               # シャード数（パスのハッシュで分散）

MAX_FILE_SIZE = 50 * 1024 * 1024   # 50MB を超えるファイルは中身を読まない
SNIPPET_LEN = 30                    # スニペットの単語数
SEARCH_LIMIT = 500                  # 一回の検索でヒットさせる最大件数
NUM_EXTRACT_WORKERS = 8             # 並列でテキスト抽出するスレッド数
BATCH_SIZE = 100                    # DBへの一括書き込みサイズ

# テキストとして読むファイル拡張子
TEXT_EXT = {
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".log",
    ".json", ".jsonl", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".html", ".htm", ".css", ".scss", ".sass", ".less",
    ".js", ".mjs", ".cjs", ".ts", ".tsx", ".jsx", ".vue", ".svelte",
    ".py", ".pyw", ".rb", ".php", ".pl", ".lua",
    ".c", ".cpp", ".cxx", ".cc", ".h", ".hpp", ".hxx",
    ".java", ".kt", ".scala", ".groovy", ".cs", ".fs", ".vb",
    ".go", ".rs", ".swift", ".m", ".mm",
    ".sql", ".psql",
    ".sh", ".bash", ".zsh", ".bat", ".cmd", ".ps1",
    ".r", ".jl", ".dart", ".tex", ".bib",
    ".gitignore", ".dockerfile", ".editorconfig",
}
OFFICE_EXT = {".docx", ".xlsx", ".pptx", ".pdf", ".doc", ".xls", ".ppt", ".rtf"}

# 走査対象のグループ分け（UIのチェックボックスと連動）
FILE_GROUPS: dict[str, set[str]] = {
    "Word":       {".docx", ".doc", ".rtf"},
    "Excel":      {".xlsx", ".xls"},
    "PowerPoint": {".pptx", ".ppt"},
    "PDF":        {".pdf"},
    "テキスト/データ": {
        ".txt", ".md", ".markdown", ".rst", ".log",
        ".csv", ".tsv", ".json", ".jsonl", ".xml",
        ".html", ".htm",
    },
    "ソースコード/設定": {
        ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
        ".css", ".scss", ".sass", ".less",
        ".js", ".mjs", ".cjs", ".ts", ".tsx", ".jsx", ".vue", ".svelte",
        ".py", ".pyw", ".rb", ".php", ".pl", ".lua",
        ".c", ".cpp", ".cxx", ".cc", ".h", ".hpp", ".hxx",
        ".java", ".kt", ".scala", ".groovy",
        ".cs", ".fs", ".vb",
        ".go", ".rs", ".swift", ".m", ".mm",
        ".sql", ".psql",
        ".sh", ".bash", ".zsh", ".bat", ".cmd", ".ps1",
        ".r", ".jl", ".dart", ".tex", ".bib",
        ".gitignore", ".dockerfile", ".editorconfig",
    },
}

# 走査をスキップするディレクトリ
SKIP_DIRS = {
    "node_modules", ".git", ".svn", ".hg", "__pycache__",
    ".venv", "venv", "env", ".env",
    ".idea", ".vscode", ".vs",
    "dist", "build", "out", "target", ".next", ".nuxt", ".cache",
    "System Volume Information", "$RECYCLE.BIN", "$Recycle.Bin",
    "AppData",
}

# ---------------------------------------------------------------------------
# テキスト抽出
# ---------------------------------------------------------------------------

def extract_text(path: Path) -> str | None:
    """各種ファイルからテキストを取り出す。読めなければ None。"""
    try:
        size = path.stat().st_size
    except OSError:
        return None
    if size > MAX_FILE_SIZE:
        return None

    ext = path.suffix.lower()

    # ドットなしファイル名（Dockerfile など）も拾う
    if ext in TEXT_EXT or path.name in {"Dockerfile", "Makefile", "Procfile", ".env"}:
        return _read_text_file(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".doc":
        return _extract_doc(path)
    if ext == ".xls":
        return _extract_xls(path)
    if ext == ".ppt":
        return _extract_ppt(path)
    if ext == ".rtf":
        return _extract_rtf(path)
    return None


def _read_text_file(path: Path) -> str | None:
    # 日本語環境のファイルを考慮した順番
    for enc in ("utf-8-sig", "utf-8", "cp932", "shift_jis", "euc-jp", "utf-16"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
        except OSError:
            return None
    # 最終手段：壊れた文字は無視
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except OSError:
        return None


def _extract_pdf(path: Path) -> str | None:
    try:
        import pypdf
    except ImportError:
        return None
    try:
        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            if reader.is_encrypted:
                try:
                    reader.decrypt("")
                except Exception:
                    return None
            parts = []
            for page in reader.pages:
                try:
                    parts.append(page.extract_text() or "")
                except Exception:
                    continue
            return "\n".join(parts)
    except Exception:
        return None


def _extract_docx(path: Path) -> str | None:
    try:
        import docx
    except ImportError:
        return None
    try:
        doc = docx.Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        parts.append(cell.text)
        return "\n".join(parts)
    except Exception:
        return None


def _extract_xlsx(path: Path) -> str | None:
    try:
        import openpyxl
    except ImportError:
        return None
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(parts)
    except Exception:
        return None


def _extract_pptx(path: Path) -> str | None:
    try:
        from pptx import Presentation
    except ImportError:
        return None
    try:
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    except Exception:
        return None


# ---- 旧Office (Office 97-2003 バイナリ形式) ----

# Word/PowerPoint COMの呼び出しを直列化（Officeアプリは同時アクセスに弱い）
_office_lock = threading.Lock()
_com_availability: dict[str, bool] = {}  # "Word.Application" → True/False のキャッシュ


def _check_com_app(prog_id: str) -> bool:
    """指定したCOMアプリが起動できるか確認（結果はキャッシュ）。"""
    if prog_id in _com_availability:
        return _com_availability[prog_id]
    if sys.platform != "win32":
        _com_availability[prog_id] = False
        return False
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        _com_availability[prog_id] = False
        return False
    pythoncom.CoInitialize()
    app = None
    ok = False
    try:
        app = win32com.client.DispatchEx(prog_id)
        ok = True
    except Exception:
        ok = False
    finally:
        try:
            if app is not None:
                app.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
    _com_availability[prog_id] = ok
    return ok


def _extract_via_word_com(path: Path) -> str | None:
    """Microsoft Word 経由で .doc を読む。要 pywin32 + Word インストール。"""
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return None

    with _office_lock:
        pythoncom.CoInitialize()
        word = None
        doc = None
        try:
            word = win32com.client.DispatchEx("Word.Application")
            word.Visible = False
            word.DisplayAlerts = 0
            doc = word.Documents.Open(
                str(path.absolute()),
                ReadOnly=True,
                AddToRecentFiles=False,
                ConfirmConversions=False,
                NoEncodingDialog=True,
            )
            parts = []
            # 本文
            try:
                parts.append(doc.Content.Text or "")
            except Exception:
                pass
            # ヘッダ／フッタ
            try:
                for section in doc.Sections:
                    for hf in (section.Headers, section.Footers):
                        for h in hf:
                            t = h.Range.Text
                            if t:
                                parts.append(t)
            except Exception:
                pass
            return "\n".join(parts)
        except Exception:
            return None
        finally:
            try:
                if doc is not None:
                    doc.Close(SaveChanges=False)
            except Exception:
                pass
            try:
                if word is not None:
                    word.Quit()
            except Exception:
                pass
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _extract_via_powerpoint_com(path: Path) -> str | None:
    """Microsoft PowerPoint 経由で .ppt を読む。要 pywin32 + PowerPoint インストール。"""
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return None

    with _office_lock:
        pythoncom.CoInitialize()
        app = None
        pres = None
        try:
            app = win32com.client.DispatchEx("PowerPoint.Application")
            # PowerPoint は Visible が変更不可な版もあるので try
            try:
                app.DisplayAlerts = 0
            except Exception:
                pass
            pres = app.Presentations.Open(
                str(path.absolute()),
                ReadOnly=True,
                Untitled=False,
                WithWindow=False,
            )
            parts = []
            for i, slide in enumerate(pres.Slides, 1):
                parts.append(f"[Slide {i}]")
                try:
                    for shape in slide.Shapes:
                        if shape.HasTextFrame and shape.TextFrame.HasText:
                            parts.append(shape.TextFrame.TextRange.Text or "")
                except Exception:
                    continue
                # ノート
                try:
                    notes = slide.NotesPage.Shapes.Placeholders(2).TextFrame.TextRange.Text
                    if notes:
                        parts.append(f"[Notes {i}] {notes}")
                except Exception:
                    pass
            return "\n".join(parts)
        except Exception:
            return None
        finally:
            try:
                if pres is not None:
                    pres.Close()
            except Exception:
                pass
            try:
                if app is not None:
                    app.Quit()
            except Exception:
                pass
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _extract_xls(path: Path) -> str | None:
    """Excel 97-2003 (.xls) - xlrd で読む。完全対応。"""
    try:
        import xlrd
    except ImportError:
        return None
    try:
        wb = xlrd.open_workbook(str(path), on_demand=True)
        parts = []
        for name in wb.sheet_names():
            sheet = wb.sheet_by_name(name)
            parts.append(f"[Sheet: {name}]")
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                cells = [str(c) for c in row if c not in (None, "")]
                if cells:
                    parts.append(" | ".join(cells))
            wb.unload_sheet(name)
        return "\n".join(parts)
    except Exception:
        return None


# 印字可能ASCII + 全角記号 + ひらがな + カタカナ + 漢字 + 半角ｶﾅ + 改行類
_READABLE_RE = re.compile(
    r"[\x20-\x7e\u3000-\u30ff\u4e00-\u9fff\uff00-\uffef\n\r\t]{4,}"
)


def _extract_ole_text(path: Path, stream_names: tuple[str, ...]) -> str | None:
    """OLEファイル (.doc/.ppt) からテキスト抽出。書式コード混在のため
    ベストエフォート（本文の大半は取れるが、見出しやヘッダの一部が失われることあり）。"""
    try:
        import olefile
    except ImportError:
        return None
    try:
        if not olefile.isOleFile(str(path)):
            return None
        with olefile.OleFileIO(str(path)) as ole:
            for sname in stream_names:
                if not ole.exists(sname):
                    continue
                data = ole.openstream(sname).read()
                # UTF-16LE と CP932 の両方を試して、より長い「読める」テキストを採用
                pieces_u16 = _READABLE_RE.findall(
                    data.decode("utf-16-le", errors="ignore")
                )
                pieces_cp932 = _READABLE_RE.findall(
                    data.decode("cp932", errors="ignore")
                )
                u16_text = "\n".join(pieces_u16)
                cp_text = "\n".join(pieces_cp932)
                # 日本語文字（ひらがな/カタカナ/漢字）の出現数で判定
                jp_re = re.compile(r"[\u3040-\u30ff\u4e00-\u9fff]")
                u16_score = len(jp_re.findall(u16_text))
                cp_score = len(jp_re.findall(cp_text))
                # 日本語が多く取れた方を採用、どちらも少ないなら長い方
                if u16_score > cp_score:
                    return u16_text
                if cp_score > u16_score:
                    return cp_text
                return u16_text if len(u16_text) >= len(cp_text) else cp_text
        return None
    except Exception:
        return None


def _extract_doc(path: Path) -> str | None:
    """Word 97-2003 (.doc) - Word(COM)が使えればそちら、ダメなら olefile ヒューリスティック。"""
    # COMが使える環境なら高品質抽出
    if _check_com_app("Word.Application"):
        text = _extract_via_word_com(path)
        if text:
            return text
    # フォールバック：純Pythonのベストエフォート
    return _extract_ole_text(path, ("WordDocument",))


def _extract_ppt(path: Path) -> str | None:
    """PowerPoint 97-2003 (.ppt) - PowerPoint(COM)優先、ダメなら olefile ヒューリスティック。"""
    if _check_com_app("PowerPoint.Application"):
        text = _extract_via_powerpoint_com(path)
        if text:
            return text
    return _extract_ole_text(path, ("PowerPoint Document",))


def _extract_rtf(path: Path) -> str | None:
    """RTF - striprtf で読む。完全対応。"""
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        return None
    try:
        # RTF はASCII互換のテキストフォーマット
        for enc in ("utf-8", "cp932", "latin-1"):
            try:
                with open(path, "r", encoding=enc) as f:
                    rtf = f.read()
                return rtf_to_text(rtf, errors="ignore")
            except UnicodeDecodeError:
                continue
        return None
    except Exception:
        return None


# ---------------------------------------------------------------------------
# インデックス（SQLite FTS5）
# ---------------------------------------------------------------------------

class Indexer:
    """シャード化されたインデックス。NUM_SHARDS 個のDBにパスのハッシュで分散保存。

    - 書き込み: パスから shard_for() でシャードを決定 → 各シャードへバッチ書き込み
    - 検索: 全シャードを並列に検索してマージ
    """

    def __init__(self, shard_dir: Path, num_shards: int = NUM_SHARDS):
        self.shard_dir = shard_dir
        self.num_shards = num_shards
        self.shard_dir.mkdir(parents=True, exist_ok=True)
        self.tokenizer = "trigram"
        self._init_all_shards()

    def shard_path(self, idx: int) -> Path:
        return self.shard_dir / f"index_{idx:02x}.db"

    def shard_for(self, path: str) -> int:
        """パスからシャード番号を決定。Windowsの大文字小文字非区別に対応。"""
        # os.path.normcase: Windowsでは小文字化＋区切り正規化、他OSではnoop
        norm = os.path.normcase(path)
        h = hashlib.md5(norm.encode("utf-8", errors="replace")).digest()
        return int.from_bytes(h[:4], "little") % self.num_shards

    def _connect(self, shard_idx: int) -> sqlite3.Connection:
        conn = sqlite3.connect(str(self.shard_path(shard_idx)), timeout=30)
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("PRAGMA synchronous=NORMAL")
        conn.execute("PRAGMA temp_store=MEMORY")
        return conn

    def _init_all_shards(self) -> None:
        for i in range(self.num_shards):
            conn = self._connect(i)
            try:
                try:
                    conn.execute(
                        "CREATE VIRTUAL TABLE IF NOT EXISTS contents USING fts5("
                        "path UNINDEXED, name, content, tokenize='trigram')"
                    )
                except sqlite3.OperationalError:
                    conn.execute(
                        "CREATE VIRTUAL TABLE IF NOT EXISTS contents USING fts5("
                        "path UNINDEXED, name, content, "
                        "tokenize='unicode61 remove_diacritics 2')"
                    )
                    self.tokenizer = "unicode61"
                conn.executescript("""
                    CREATE TABLE IF NOT EXISTS files (
                        path      TEXT PRIMARY KEY,
                        name      TEXT,
                        mtime     REAL,
                        size      INTEGER,
                        indexed   REAL
                    );
                    CREATE INDEX IF NOT EXISTS idx_files_name ON files(name);
                """)
                conn.commit()
            finally:
                conn.close()

    # --- 書き込み ---
    def get_known(self) -> dict[str, float]:
        """全シャードを横断して {path: mtime} を返す。差分判定用。"""
        result: dict[str, float] = {}
        for i in range(self.num_shards):
            conn = self._connect(i)
            try:
                for p, m in conn.execute("SELECT path, mtime FROM files"):
                    result[p] = m
            finally:
                conn.close()
        return result

    def upsert_batch(self, items: list[tuple[str, str, float, int, str]]) -> None:
        """items を該当シャードに振り分けて一括書き込み。"""
        if not items:
            return
        by_shard: dict[int, list] = {}
        for item in items:
            by_shard.setdefault(self.shard_for(item[0]), []).append(item)
        for sidx, shard_items in by_shard.items():
            conn = self._connect(sidx)
            try:
                with conn:
                    paths = [(p,) for p, *_ in shard_items]
                    conn.executemany("DELETE FROM contents WHERE path=?", paths)
                    conn.executemany(
                        "INSERT INTO contents(path, name, content) VALUES(?,?,?)",
                        [(p, n, c or "") for p, n, _, _, c in shard_items],
                    )
                    now = time.time()
                    conn.executemany(
                        "INSERT OR REPLACE INTO files(path, name, mtime, size, indexed) "
                        "VALUES(?,?,?,?,?)",
                        [(p, n, m, s, now) for p, n, m, s, _ in shard_items],
                    )
            finally:
                conn.close()

    def delete_paths(self, paths: list[str]) -> None:
        if not paths:
            return
        by_shard: dict[int, list] = {}
        for p in paths:
            by_shard.setdefault(self.shard_for(p), []).append(p)
        for sidx, sp in by_shard.items():
            conn = self._connect(sidx)
            try:
                with conn:
                    conn.executemany("DELETE FROM contents WHERE path=?", [(p,) for p in sp])
                    conn.executemany("DELETE FROM files WHERE path=?", [(p,) for p in sp])
            finally:
                conn.close()

    def clear(self) -> None:
        for i in range(self.num_shards):
            conn = self._connect(i)
            try:
                with conn:
                    conn.execute("DELETE FROM contents")
                    conn.execute("DELETE FROM files")
            finally:
                conn.close()

    def stats(self) -> tuple[int, int]:
        total_n = 0
        total_size = 0
        for i in range(self.num_shards):
            conn = self._connect(i)
            try:
                row = conn.execute(
                    "SELECT COUNT(*), COALESCE(SUM(size),0) FROM files"
                ).fetchone()
                total_n += row[0]
                total_size += row[1]
            finally:
                conn.close()
        return total_n, total_size

    def per_shard_stats(self) -> list[tuple[int, int, int]]:
        """[(shard_idx, file_count, db_bytes), ...] を返す。"""
        result = []
        for i in range(self.num_shards):
            conn = self._connect(i)
            try:
                n = conn.execute("SELECT COUNT(*) FROM files").fetchone()[0]
            finally:
                conn.close()
            db_bytes = 0
            for suffix in ("", "-wal", "-shm"):
                p = Path(str(self.shard_path(i)) + suffix)
                if p.exists():
                    try:
                        db_bytes += p.stat().st_size
                    except OSError:
                        pass
            result.append((i, n, db_bytes))
        return result

    def get_contents(self, paths: list[str]) -> dict[str, str]:
        """複数パスの本文を一括取得。"""
        if not paths:
            return {}
        # シャードごとにグルーピングして一発クエリ
        by_shard: dict[int, list[str]] = {}
        for p in paths:
            by_shard.setdefault(self.shard_for(p), []).append(p)
        result: dict[str, str] = {}
        for sidx, sps in by_shard.items():
            conn = self._connect(sidx)
            try:
                # IN クエリの長さ制限を考慮してチャンク
                for i in range(0, len(sps), 500):
                    chunk = sps[i:i + 500]
                    placeholders = ",".join("?" for _ in chunk)
                    for p, c in conn.execute(
                        f"SELECT path, content FROM contents WHERE path IN ({placeholders})",
                        chunk,
                    ):
                        result[p] = c or ""
            finally:
                conn.close()
        return result

    # --- 検索（全シャード並列） ---
    def search(self, query: str, mode: str = "AND",
               limit: int = SEARCH_LIMIT) -> list[tuple[str, str, str]]:
        if not query.strip():
            return []
        mode = "OR" if mode.upper() == "OR" else "AND"

        tokens = _parse_query_tokens(query)
        fts_pos: list[tuple[str, bool]] = []
        fts_neg: list[tuple[str, bool]] = []
        like_pos: list[str] = []
        like_neg: list[str] = []
        for term, is_phrase, is_neg in tokens:
            if _len_no_space(term) >= 3:
                (fts_neg if is_neg else fts_pos).append((term, is_phrase))
            else:
                (like_neg if is_neg else like_pos).append(term)

        if not fts_pos and not like_pos:
            return []

        # 各シャードから少し多めに取得して、後でグローバルソート
        per_shard = max(50, (limit * 3) // self.num_shards + 20)

        all_results: list[tuple[float, str, str, str]] = []  # (rank, path, name, snip)

        with ThreadPoolExecutor(
            max_workers=min(self.num_shards, 8),
            thread_name_prefix="search-shard"
        ) as pool:
            futures = [
                pool.submit(self._search_one_shard, i,
                            fts_pos, fts_neg, like_pos, like_neg, mode, per_shard)
                for i in range(self.num_shards)
            ]
            for f in futures:
                try:
                    all_results.extend(f.result())
                except Exception:
                    continue

        # rank 昇順（小さいほど関連度高）、rank が無いLIKEヒットは最後に
        all_results.sort(key=lambda r: r[0])
        return [(p, n, s) for _, p, n, s in all_results[:limit]]

    def _search_one_shard(
        self, shard_idx: int,
        fts_pos: list[tuple[str, bool]], fts_neg: list[tuple[str, bool]],
        like_pos: list[str], like_neg: list[str],
        mode: str, limit: int,
    ) -> list[tuple[float, str, str, str]]:
        """1シャードの検索。戻り値は (rank, path, name, snippet) のリスト。"""
        conn = self._connect(shard_idx)
        try:
            candidates: list[tuple[float, str, str, str]] = []
            seen_paths: set[str] = set()

            # Step 1: FTS
            if fts_pos:
                op = " OR " if mode == "OR" else " AND "
                fts_expr = op.join(_fts_phrase(t) for t, _ in fts_pos)
                for t, _ in fts_neg:
                    fts_expr = f"({fts_expr}) NOT {_fts_phrase(t)}"
                try:
                    rows = conn.execute(
                        f"""
                        SELECT rank, path, name,
                               snippet(contents, 2, '〘', '〙', '…', {SNIPPET_LEN})
                        FROM contents
                        WHERE contents MATCH ?
                        ORDER BY rank
                        LIMIT ?
                        """,
                        (fts_expr, limit * 3),
                    ).fetchall()
                    for rank, path, name, snip in rows:
                        candidates.append((rank, path, name, snip))
                        seen_paths.add(path)
                except sqlite3.OperationalError:
                    pass

            # Step 2: OR + 短い肯定語 → LIKE で追加
            if mode == "OR" and like_pos:
                like_clauses = " OR ".join(
                    "content LIKE ? OR name LIKE ?" for _ in like_pos
                )
                params = []
                for t in like_pos:
                    params.extend([f"%{t}%", f"%{t}%"])
                params.append(limit * 3)
                like_rows = conn.execute(
                    f"SELECT path, name, content FROM contents WHERE {like_clauses} LIMIT ?",
                    params,
                ).fetchall()
                for path, name, content in like_rows:
                    if path in seen_paths:
                        continue
                    body = content or ""
                    snip = ""
                    for t in like_pos:
                        if t.lower() in body.lower():
                            snip = _make_snippet(body, t)
                            break
                        if t.lower() in name.lower():
                            snip = name
                            break
                    # LIKEヒットはFTS rankが無いので大きめの値（後ろに来る）
                    candidates.append((1e6, path, name, snip))
                    seen_paths.add(path)

            # Step 3: AND モードでFTS語なし → 全件走査
            if mode == "AND" and not fts_pos and like_pos:
                for path, name, content in conn.execute(
                    "SELECT path, name, content FROM contents"
                ):
                    body = content or ""
                    haystack = (body + "\n" + name).lower()
                    if all(t.lower() in haystack for t in like_pos):
                        snip = _make_snippet(body, like_pos[0])
                        candidates.append((1e6, path, name, snip))
                        seen_paths.add(path)

            # Step 4: 後処理フィルタ
            need_post_filter = (
                (mode == "AND" and like_pos and fts_pos) or like_neg
            )
            if need_post_filter and candidates:
                paths = [c[1] for c in candidates]
                contents_map: dict[str, str] = {}
                placeholders = ",".join("?" for _ in paths)
                for p, c in conn.execute(
                    f"SELECT path, content FROM contents WHERE path IN ({placeholders})",
                    paths,
                ):
                    contents_map[p] = c or ""

                filtered = []
                for rank, path, name, snip in candidates:
                    body = contents_map.get(path, "")
                    haystack = (body + "\n" + name).lower()
                    if mode == "AND" and like_pos and fts_pos:
                        if not all(t.lower() in haystack for t in like_pos):
                            continue
                    if any(t.lower() in haystack for t in like_neg):
                        continue
                    filtered.append((rank, path, name, snip))
                    if len(filtered) >= limit:
                        break
                candidates = filtered

            return candidates[:limit]
        finally:
            conn.close()


# ---------------------------------------------------------------------------
# 旧単一DBから分割DBへのマイグレーション
# ---------------------------------------------------------------------------

def migrate_from_legacy(legacy_db: Path, indexer: Indexer,
                        progress_cb=None, batch_size: int = 500) -> int:
    """旧 index.db からシャードDBへ全レコードをコピー。完了件数を返す。

    実装ノート: contents は FTS5 で path カラムが UNINDEXED のため、
    files との JOIN が O(N²) で実質終わらない。代わりに:
      1) files のメタ情報を全部メモリへロード（dict）
      2) contents を順次スキャンしながら dict で合体
      3) contents に無い files 行（ファイル名のみ）を最後に追加
    これで O(N) で済む。
    """
    if not legacy_db.exists():
        return 0
    src = sqlite3.connect(str(legacy_db))
    try:
        # 早めに progress を出して「準備中」状態を抜ける
        if progress_cb:
            progress_cb(0, 0)

        # 総件数（COUNT(*) は数秒〜十数秒）
        try:
            total = src.execute("SELECT COUNT(*) FROM files").fetchone()[0]
        except sqlite3.OperationalError:
            total = 0
        if total == 0:
            try:
                total = src.execute("SELECT COUNT(*) FROM contents").fetchone()[0]
            except sqlite3.OperationalError:
                total = 0
        if total == 0:
            return 0
        if progress_cb:
            progress_cb(0, total)

        # Phase 1: files のメタ情報を全部メモリへ
        # 100万件でも数百MB程度なのでメモリに乗る
        files_meta: dict[str, tuple[str, float, int]] = {}
        try:
            for path, name, mtime, size in src.execute(
                "SELECT path, name, mtime, size FROM files"
            ):
                files_meta[path] = (name, mtime, size)
        except sqlite3.OperationalError:
            pass  # files テーブルが無くても続行

        batch: list = []
        done = 0
        seen_in_contents: set[str] = set()

        # Phase 2: contents を順次スキャン（FTS5 の単純な全件走査は高速）
        try:
            for path, name, content in src.execute(
                "SELECT path, name, content FROM contents"
            ):
                seen_in_contents.add(path)
                meta = files_meta.get(path)
                if meta:
                    name_meta, mtime, size = meta
                    if name_meta:
                        name = name_meta
                else:
                    mtime, size = 0.0, 0
                batch.append((path, name, mtime, size, content or ""))
                if len(batch) >= batch_size:
                    indexer.upsert_batch(batch)
                    done += len(batch)
                    batch.clear()
                    if progress_cb:
                        progress_cb(done, total)
        except sqlite3.OperationalError:
            pass  # contents が無くても続行

        # Phase 3: contents に無く files にだけある行
        for path, (name, mtime, size) in files_meta.items():
            if path in seen_in_contents:
                continue
            batch.append((path, name, mtime, size, ""))
            if len(batch) >= batch_size:
                indexer.upsert_batch(batch)
                done += len(batch)
                batch.clear()
                if progress_cb:
                    progress_cb(done, total)

        if batch:
            indexer.upsert_batch(batch)
            done += len(batch)
            if progress_cb:
                progress_cb(done, total)

        return done
    finally:
        src.close()


def _len_no_space(s: str) -> int:
    return len(re.sub(r"\s", "", s))


def _parse_query_tokens(q: str) -> list[tuple[str, bool, bool]]:
    """[(term, is_phrase, is_negated), ...] を返す。"""
    out = []
    for tok in re.findall(r'-?"[^"]+"|-?\S+', q):
        neg = tok.startswith("-")
        if neg:
            tok = tok[1:]
        if tok.startswith('"') and tok.endswith('"') and len(tok) >= 2:
            out.append((tok[1:-1], True, neg))
        else:
            out.append((tok, False, neg))
    return out


def _fts_phrase(term: str) -> str:
    safe = term.replace('"', '""')
    return f'"{safe}"'


def _make_snippet(body: str, term: str, ctx: int = 60) -> str:
    if not body or not term:
        return body[:160] if body else ""
    i = body.lower().find(term.lower())
    if i < 0:
        return body[:160]
    s = max(0, i - ctx)
    e = min(len(body), i + len(term) + ctx)
    before = body[s:i].replace("\n", " ")
    hit = body[i:i + len(term)]
    after = body[i + len(term):e].replace("\n", " ")
    pre = "…" if s > 0 else ""
    suf = "…" if e < len(body) else ""
    return f"{pre}{before}〘{hit}〙{after}{suf}"


def _build_ai_markdown(
    query: str,
    instruction: str,
    rows: list[tuple[str, str, str]],
    contents_map: dict[str, str],
    max_chars_per_file: int,
) -> str:
    """検索結果をAI向けのMarkdownに整形。"""
    parts = []
    parts.append("# 検索クエリ\n")
    parts.append(f"\n```\n{query}\n```\n")

    if instruction:
        parts.append("\n# AIへの指示\n")
        parts.append(f"\n{instruction}\n")

    parts.append(f"\n# 検索結果（{len(rows)}件）\n")
    parts.append(
        "\n以下に検索でヒットしたファイルを列挙します。"
        "重要なファイルに言及する際は必ずパスを引用してください。\n"
    )

    for i, (path, name, snip) in enumerate(rows, 1):
        body = contents_map.get(path, "") or ""
        truncated = False
        if len(body) > max_chars_per_file:
            body = body[:max_chars_per_file]
            truncated = True

        parts.append("\n---\n")
        parts.append(f"\n## [{i}/{len(rows)}] {name}\n")
        parts.append(f"\n- **パス**: `{path}`\n")
        if snip:
            clean_snip = snip.replace("〘", "**").replace("〙", "**")
            parts.append(f"- **ヒット箇所**: {clean_snip}\n")
        parts.append("\n### 本文\n")
        if body.strip():
            parts.append(f"\n```\n{body}\n")
            if truncated:
                parts.append(
                    f"\n…（先頭 {max_chars_per_file:,} 文字のみ。"
                    f"実ファイルはより長い）…\n"
                )
            parts.append("```\n")
        else:
            parts.append("\n（本文取得不可：ファイル名のみインデックス）\n")

    parts.append("\n---\n\n以上です。指示に従って回答してください。\n")
    return "".join(parts)


# ---------------------------------------------------------------------------
# 走査ワーカー
# ---------------------------------------------------------------------------

def _extract_task(path_str: str, name: str, mtime: float, size: int):
    """ThreadPoolExecutor 用ワーカー。
    戻り値: (path, name, mtime, size, text, read_error)
      read_error: True = 中身を読もうとして失敗。False = 成功 or そもそも対応外拡張子。
    """
    text = None
    read_error = False
    try:
        text = extract_text(Path(path_str))
    except Exception:
        read_error = True
        text = None
    # extract_text が None を返すケース：
    #  - 対応外拡張子 → read_error=False
    #  - ファイル読み込みエラー → read_error=True
    # 区別がつかないので、対応している拡張子なら None=エラー扱いにする
    if text is None and not read_error:
        ext = Path(path_str).suffix.lower()
        # この拡張子は本来読めるはずのもの → 読み取り失敗扱い
        if ext in {".pdf", ".docx", ".xlsx", ".pptx",
                    ".doc", ".xls", ".ppt", ".rtf"} or ext in TEXT_EXT:
            read_error = True
    return (path_str, name, mtime, size, text or "", read_error)


class IndexWorker(threading.Thread):
    """バックグラウンドでフォルダを走査してインデックス化。
    抽出は ThreadPoolExecutor で並列実行（I/Oバウンドなので効果大）。
    allowed_exts が None なら全拡張子、setなら該当拡張子のみ。
    """

    def __init__(self, roots: list[str], indexer: Indexer, progress_q: queue.Queue,
                 full_rescan: bool = False,
                 allowed_exts: set[str] | None = None):
        super().__init__(daemon=True)
        self.roots = [Path(r) for r in roots]
        self.indexer = indexer
        self.q = progress_q
        self.full_rescan = full_rescan
        self.allowed_exts = allowed_exts
        self.stop_flag = threading.Event()

    def stop(self):
        self.stop_flag.set()

    def run(self):
        try:
            self._run()
        except Exception:
            self.q.put(("error", traceback.format_exc()))
        finally:
            self.q.put(("done", None))

    def _run(self):
        known = {} if self.full_rescan else self.indexer.get_known()
        seen: set[str] = set()
        batch: list[tuple[str, str, float, int, str]] = []

        scanned = 0
        indexed = 0
        skipped = 0
        read_errors = 0
        access_errors = 0
        t0 = time.time()
        last_update = 0.0
        last_dir = ""
        current_file = ""

        def push_status(force=False):
            nonlocal last_update
            now = time.time()
            if not force and now - last_update < 0.3:
                return
            last_update = now
            cf = f"  📄 {current_file}" if current_file else ""
            self.q.put(("status",
                f"走査中: {scanned:,} 件 / 取込 {indexed:,} 件 "
                f"({now - t0:.1f}s)  📂 {last_dir}{cf}"))

        def walk_err(err):
            nonlocal access_errors
            access_errors += 1
            tgt = getattr(err, "filename", None) or str(err)
            self.q.put(("log", f"⚠ アクセス不可: {tgt}"))

        def collect(fut):
            """完了したFutureから結果を取り出してバッチに追加。"""
            nonlocal indexed, read_errors
            try:
                res = fut.result()
            except Exception:
                return
            if res is None:
                return
            path_str, name, mtime, size, text, read_error = res
            if read_error:
                read_errors += 1
                if read_errors <= 20:  # ログが溢れないよう先頭だけ
                    self.q.put(("log", f"⚠ 読み込み失敗: {path_str}"))
                elif read_errors == 21:
                    self.q.put(("log", "⚠ 読み込み失敗が21件目。以降はサマリのみ。"))
            batch.append((path_str, name, mtime, size, text))
            indexed += 1
            if len(batch) >= BATCH_SIZE:
                self.indexer.upsert_batch(batch)
                batch.clear()
                push_status(force=True)

        # 並列で抽出 - I/Oバウンドなのでスレッド数を多めに
        max_inflight = NUM_EXTRACT_WORKERS * 4

        with ThreadPoolExecutor(max_workers=NUM_EXTRACT_WORKERS,
                                thread_name_prefix="extract") as pool:
            inflight: list = []

            def shutdown():
                """停止要求時の片付け。"""
                pool.shutdown(wait=False, cancel_futures=True)

            for root in self.roots:
                if not root.exists():
                    self.q.put(("log", f"⚠ パスが見つからない: {root}"))
                    continue
                self.q.put(("log", f"スキャン開始: {root}"))
                push_status(force=True)

                for dirpath, dirnames, filenames in os.walk(
                    root, topdown=True, onerror=walk_err
                ):
                    if self.stop_flag.is_set():
                        shutdown()
                        return
                    dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS]
                    last_dir = dirpath
                    current_file = ""
                    push_status()

                    for fn in filenames:
                        if self.stop_flag.is_set():
                            shutdown()
                            return
                        # 拡張子フィルタ
                        if self.allowed_exts is not None:
                            ext = Path(fn).suffix.lower()
                            if ext not in self.allowed_exts:
                                continue
                        p = Path(dirpath) / fn
                        sp = str(p)
                        seen.add(sp)
                        scanned += 1
                        current_file = fn
                        push_status()

                        try:
                            st = p.stat()
                        except OSError as e:
                            skipped += 1
                            if skipped <= 20:
                                self.q.put(("log",
                                    f"⚠ stat失敗（スキップ）: {sp} ({e})"))
                            elif skipped == 21:
                                self.q.put(("log",
                                    "⚠ stat失敗が21件目。以降はサマリのみ。"))
                            continue

                        # 増分判定
                        prev = known.get(sp)
                        if prev is not None and abs(prev - st.st_mtime) < 0.001:
                            continue

                        # バックプレッシャ：保留中が多すぎたら少し待つ
                        if len(inflight) >= max_inflight:
                            done, _ = wait(inflight, return_when=FIRST_COMPLETED)
                            for f in done:
                                collect(f)
                                inflight.remove(f)

                        # 抽出をプールに投げる
                        fut = pool.submit(_extract_task, sp, fn,
                                          st.st_mtime, st.st_size)
                        inflight.append(fut)

            # 全投入完了 → 残りを回収
            while inflight:
                if self.stop_flag.is_set():
                    shutdown()
                    return
                done, _ = wait(inflight, return_when=FIRST_COMPLETED)
                for f in done:
                    collect(f)
                    inflight.remove(f)
                push_status()

        if batch:
            self.indexer.upsert_batch(batch)

        # 既知だが今回見つからなかったファイルをDBから削除
        if not self.full_rescan:
            roots_str = [str(r) for r in self.roots]
            stale = [p for p in known
                     if any(p == r or p.startswith(r + os.sep) for r in roots_str)
                     and p not in seen]
            if stale:
                self.indexer.delete_paths(stale)
                self.q.put(("log", f"削除済みファイルをインデックスから除去: {len(stale)} 件"))

        elapsed = time.time() - t0
        summary = (f"完了: 走査 {scanned:,} 件 / "
                   f"取込 {indexed:,} 件 / "
                   f"stat失敗 {skipped:,} / "
                   f"読込失敗 {read_errors:,} / "
                   f"アクセス不可フォルダ {access_errors:,} / "
                   f"{elapsed:.1f}s")
        self.q.put(("log", summary))


# ---------------------------------------------------------------------------
# 設定保存
# ---------------------------------------------------------------------------

def load_config() -> dict:
    if CONFIG_PATH.exists():
        try:
            return json.loads(CONFIG_PATH.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"roots": []}


def save_config(cfg: dict) -> None:
    CONFIG_PATH.write_text(json.dumps(cfg, ensure_ascii=False, indent=2), encoding="utf-8")


# ---------------------------------------------------------------------------
# GUI
# ---------------------------------------------------------------------------

class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title(f"{APP_NAME} — 全文検索")

        # WindowsのDPIスケーリングを正しく扱う（Tkがぼやけるのを防ぐ）
        if sys.platform == "win32":
            try:
                from ctypes import windll
                # Win 8.1+：モニタごとのDPI対応
                windll.shcore.SetProcessDpiAwareness(1)
            except Exception:
                try:
                    windll.user32.SetProcessDPIAware()
                except Exception:
                    pass

        # ウィンドウサイズを画面サイズに合わせる（最大でも画面の85%）
        self.update_idletasks()
        sw = self.winfo_screenwidth()
        sh = self.winfo_screenheight()
        # タスクバー分を引いた大体の作業可能高さ
        usable_h = sh - 80
        w = min(1200, int(sw * 0.85))
        h = min(800, int(usable_h * 0.95))
        x = max(0, (sw - w) // 2)
        y = max(0, (usable_h - h) // 2)
        self.geometry(f"{w}x{h}+{x}+{y}")
        self.minsize(720, 480)

        try:
            self.iconbitmap(default="")
        except Exception:
            pass

        # 設定からシャードDB の場所を取得（なければデフォルト）
        self.cfg = load_config()
        configured = self.cfg.get("db_dir")
        if configured and Path(configured).exists():
            shard_dir = Path(configured)
        else:
            shard_dir = SHARD_DIR
        self.indexer = Indexer(shard_dir)
        self.worker: IndexWorker | None = None
        self.q: queue.Queue = queue.Queue()
        self._last_stats_update = 0.0

        self._build_ui()
        self._refresh_stats()
        # 旧単一DBが残っていればマイグレーション提案
        self.after(200, self._maybe_migrate_legacy)
        self.after(100, self._poll_queue)

    # ---- UI ----
    def _build_ui(self):
        style = ttk.Style(self)
        try:
            style.theme_use("vista")
        except tk.TclError:
            pass

        # 上：対象フォルダ
        top = ttk.LabelFrame(self, text="検索対象フォルダ（ローカル / \\\\server\\share）")
        top.pack(fill="x", padx=8, pady=(8, 4))

        list_frame = ttk.Frame(top)
        list_frame.pack(side="left", fill="x", expand=True, padx=(8, 4), pady=8)
        self.roots_var = tk.StringVar()
        self.roots_list = tk.Listbox(list_frame, height=4, selectmode="extended",
                                      activestyle="dotbox")
        for r in self.cfg.get("roots", []):
            self.roots_list.insert("end", r)
        rsb = ttk.Scrollbar(list_frame, orient="vertical",
                             command=self.roots_list.yview)
        self.roots_list.configure(yscrollcommand=rsb.set)
        self.roots_list.pack(side="left", fill="both", expand=True)
        rsb.pack(side="right", fill="y")

        btns = ttk.Frame(top)
        btns.pack(side="right", padx=8, pady=8)
        ttk.Button(btns, text="フォルダ追加…", command=self.add_folder).pack(fill="x")
        ttk.Button(btns, text="UNCパス追加…", command=self.add_unc).pack(fill="x", pady=(4, 0))
        ttk.Button(btns, text="削除", command=self.remove_folder).pack(fill="x", pady=(4, 0))

        # 中：対象ファイル種別のフィルタ
        flt = ttk.LabelFrame(self, text="インデックス対象（チェックを外すと次回スキャン以降は無視）")
        flt.pack(fill="x", padx=8, pady=4)
        self.filter_vars: dict[str, tk.BooleanVar] = {}
        saved_filters = self.cfg.get("filters", {})

        # 2行レイアウト：1行に最大4個（6個なら2行目に2個）
        COLS = 4
        items = list(FILE_GROUPS.items())
        for i, (name, exts) in enumerate(items):
            var = tk.BooleanVar(value=saved_filters.get(name, True))
            self.filter_vars[name] = var
            label = f"{name}  ({len(exts)})"
            cb = ttk.Checkbutton(flt, text=label, variable=var,
                                 command=self._save_filters)
            cb.grid(row=i // COLS, column=i % COLS,
                    padx=10, pady=4, sticky="w")
        for c in range(COLS):
            flt.grid_columnconfigure(c, weight=1, uniform="filters")

        # 一括ON/OFFボタン
        btn_row = ttk.Frame(flt)
        btn_row.grid(row=(len(items) + COLS - 1) // COLS, column=0,
                     columnspan=COLS, sticky="w", padx=6, pady=(0, 4))
        ttk.Button(btn_row, text="全部ON",
                   command=lambda: self._set_all_filters(True),
                   width=8).pack(side="left", padx=2)
        ttk.Button(btn_row, text="全部OFF",
                   command=lambda: self._set_all_filters(False),
                   width=8).pack(side="left", padx=2)

        # 中：操作バー
        ops = ttk.Frame(self)
        ops.pack(fill="x", padx=8, pady=4)
        self.btn_index = ttk.Button(ops, text="インデックス更新（増分）", command=self.start_index)
        self.btn_index.pack(side="left")
        self.btn_full = ttk.Button(ops, text="フルスキャン", command=lambda: self.start_index(full=True))
        self.btn_full.pack(side="left", padx=(4, 0))
        self.btn_stop = ttk.Button(ops, text="停止", command=self.stop_index, state="disabled")
        self.btn_stop.pack(side="left", padx=(4, 0))
        ttk.Button(ops, text="DBクリア", command=self.clear_db).pack(side="left", padx=(4, 0))
        ttk.Button(ops, text="DB設定…", command=self.open_db_settings).pack(side="left", padx=(4, 0))
        ttk.Button(ops, text="🤖 AI用テキスト出力…",
                   command=self.open_ai_export).pack(side="left", padx=(4, 0))

        # スキャン中のアニメーション
        self.progress = ttk.Progressbar(ops, mode="indeterminate", length=160)
        self.progress.pack(side="left", padx=(8, 0))

        self.stats_var = tk.StringVar(value="")
        ttk.Label(ops, textvariable=self.stats_var).pack(side="right")

        # 検索バー
        sf = ttk.LabelFrame(self, text="検索（\"フレーズ\" でフレーズ完全一致、-キーワード で除外）")
        sf.pack(fill="x", padx=8, pady=4)

        # 1行目：入力ボックスと検索ボタン
        row1 = ttk.Frame(sf)
        row1.pack(fill="x", padx=8, pady=(8, 4))
        self.query_var = tk.StringVar()
        e = ttk.Entry(row1, textvariable=self.query_var, font=("", 12))
        e.pack(side="left", fill="x", expand=True)
        e.bind("<Return>", lambda _: self.do_search())
        ttk.Button(row1, text="検索", command=self.do_search).pack(side="right", padx=(8, 0))

        # 2行目：AND/OR モード選択
        row2 = ttk.Frame(sf)
        row2.pack(fill="x", padx=8, pady=(0, 8))
        ttk.Label(row2, text="複数キーワード:").pack(side="left")
        self.search_mode = tk.StringVar(value=self.cfg.get("search_mode", "AND"))
        ttk.Radiobutton(row2, text="すべて含む (AND)", value="AND",
                        variable=self.search_mode,
                        command=self._save_search_mode).pack(side="left", padx=(8, 4))
        ttk.Radiobutton(row2, text="いずれか含む (OR)", value="OR",
                        variable=self.search_mode,
                        command=self._save_search_mode).pack(side="left", padx=(4, 0))

        # 結果ペイン（左：ファイル一覧、右：プレビュー）
        # ※ body 自体の pack はメソッド最後に移動（log と status を先に下端固定するため）
        body = ttk.Panedwindow(self, orient="horizontal")

        left = ttk.Frame(body)
        body.add(left, weight=1)
        self.results = ttk.Treeview(
            left, columns=("path",), show="tree headings", height=20
        )
        self.results.heading("#0", text="ファイル名")
        self.results.heading("path", text="パス")
        self.results.column("#0", width=240, stretch=False, minwidth=120)
        self.results.column("path", width=600, minwidth=200)

        vsb = ttk.Scrollbar(left, orient="vertical", command=self.results.yview)
        hsb = ttk.Scrollbar(left, orient="horizontal", command=self.results.xview)
        self.results.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)

        # gridで確実に縦・横スクロールバーを表示
        self.results.grid(row=0, column=0, sticky="nsew")
        vsb.grid(row=0, column=1, sticky="ns")
        hsb.grid(row=1, column=0, sticky="ew")
        left.grid_rowconfigure(0, weight=1)
        left.grid_columnconfigure(0, weight=1)

        self.results.bind("<<TreeviewSelect>>", self.on_select)
        self.results.bind("<Double-1>", self.on_open)
        # マウスホイール（Windowsだと自動だが念のため）
        self.results.bind("<MouseWheel>",
            lambda e: self.results.yview_scroll(int(-e.delta/120), "units"))

        # 右クリックメニュー
        self.ctx_menu = tk.Menu(self, tearoff=0)
        self.ctx_menu.add_command(label="ファイルを開く", command=self.on_open)
        self.ctx_menu.add_command(label="場所を開く（エクスプローラ）",
                                   command=self.open_location)
        self.ctx_menu.add_separator()
        self.ctx_menu.add_command(label="パスをコピー", command=self.copy_path)
        self.ctx_menu.add_command(label="ファイル名をコピー", command=self.copy_name)
        # Windows/Linux=Button-3, macOS=Button-2
        self.results.bind("<Button-3>", self._show_ctx_menu)
        self.results.bind("<Button-2>", self._show_ctx_menu)

        right = ttk.Frame(body)
        body.add(right, weight=1)
        ttk.Label(right, text="プレビュー（ヒット箇所）").pack(anchor="w")
        self.preview = scrolledtext.ScrolledText(right, wrap="word", font=("Consolas", 10))
        self.preview.pack(fill="both", expand=True)
        self.preview.tag_configure("hit", background="#fff3a0", foreground="#000")
        self.preview.tag_configure("path", foreground="#0066cc", font=("", 10, "bold"))

        # ===== 画面下端に固定する要素を「先に」packする（side='bottom'）=====
        # 状態バー（最下段）
        self.status_var = tk.StringVar(value="準備完了")
        ttk.Label(self, textvariable=self.status_var,
                  relief="sunken", anchor="w").pack(side="bottom", fill="x")

        # ログ（状態バーの上）— 下端固定なので body が縮んでも常に見える
        log_frame = ttk.LabelFrame(self, text="ログ")
        log_frame.pack(side="bottom", fill="x", padx=8, pady=(4, 6))
        self.log = scrolledtext.ScrolledText(log_frame, height=4, wrap="word")
        self.log.pack(fill="x", padx=4, pady=4)

        # ===== 最後に body をpack（残りスペース全部に展開） =====
        body.pack(side="top", fill="both", expand=True, padx=8, pady=4)

        self._results_data: list[tuple[str, str, str]] = []

    # ---- フォルダ操作 ----
    def add_folder(self):
        d = filedialog.askdirectory(title="検索対象フォルダを選択")
        if d:
            self.roots_list.insert("end", d)
            self._save_roots()

    def add_unc(self):
        win = tk.Toplevel(self)
        win.title("UNCパスを追加")
        win.geometry("500x120")
        win.transient(self)
        ttk.Label(win, text=r"例: \\server\share\folder").pack(pady=(10, 0))
        v = tk.StringVar()
        e = ttk.Entry(win, textvariable=v, width=60)
        e.pack(pady=8, padx=10, fill="x")
        e.focus()
        def ok():
            p = v.get().strip().strip('"')
            if p:
                self.roots_list.insert("end", p)
                self._save_roots()
                win.destroy()
        ttk.Button(win, text="追加", command=ok).pack()
        e.bind("<Return>", lambda _: ok())

    def remove_folder(self):
        for i in reversed(self.roots_list.curselection()):
            self.roots_list.delete(i)
        self._save_roots()

    def _save_roots(self):
        self.cfg["roots"] = list(self.roots_list.get(0, "end"))
        save_config(self.cfg)

    def _save_filters(self):
        self.cfg["filters"] = {n: v.get() for n, v in self.filter_vars.items()}
        save_config(self.cfg)

    def _set_all_filters(self, value: bool):
        for v in self.filter_vars.values():
            v.set(value)
        self._save_filters()

    def _save_search_mode(self):
        self.cfg["search_mode"] = self.search_mode.get()
        save_config(self.cfg)

    def _allowed_exts(self) -> set[str]:
        """チェック中のグループに属する拡張子の集合を返す。"""
        allowed: set[str] = set()
        for name, var in self.filter_vars.items():
            if var.get():
                allowed |= FILE_GROUPS[name]
        return allowed

    # ---- インデックス ----
    def start_index(self, full: bool = False):
        if self.worker and self.worker.is_alive():
            return
        roots = list(self.roots_list.get(0, "end"))
        if not roots:
            messagebox.showwarning(APP_NAME, "検索対象フォルダを1件以上追加してください。")
            return
        allowed = self._allowed_exts()
        if not allowed:
            messagebox.showwarning(APP_NAME,
                "「インデックス対象」のチェックボックスを最低1つは選んでください。")
            return
        self.btn_index.config(state="disabled")
        self.btn_full.config(state="disabled")
        self.btn_stop.config(state="normal")
        self.progress.start(10)
        groups_on = [n for n, v in self.filter_vars.items() if v.get()]
        self._log(f"--- {'フルスキャン' if full else '増分インデックス'} 開始 "
                  f"（対象: {', '.join(groups_on)}）---")
        self.worker = IndexWorker(roots, self.indexer, self.q,
                                  full_rescan=full, allowed_exts=allowed)
        self.worker.start()

    def stop_index(self):
        if self.worker and self.worker.is_alive():
            self.worker.stop()
            self._log("停止要求...")

    def clear_db(self):
        if not messagebox.askyesno(APP_NAME, "インデックスを全消去します。よろしいですか？"):
            return
        self.indexer.clear()
        self._log("インデックスをクリアしました。")
        self._refresh_stats()

    # ---- DB設定 ----
    # ---- AI用テキスト出力 ----
    def open_ai_export(self):
        """検索結果をAI（Copilot/ChatGPT/Claude等）に渡しやすい形式で出力するダイアログ。"""
        win = tk.Toplevel(self)
        win.title("AI用テキスト出力")
        win.geometry("620x520")
        win.transient(self)

        frm = ttk.Frame(win)
        frm.pack(fill="both", expand=True, padx=14, pady=12)

        # クエリ
        ttk.Label(frm, text="検索クエリ（何について知りたいか）:",
                  font=("", 10, "bold")).pack(anchor="w")
        query_var = tk.StringVar(value=self.query_var.get())
        ttk.Entry(frm, textvariable=query_var, font=("", 11)).pack(fill="x", pady=(2, 8))

        # AIへの指示
        ttk.Label(frm, text="AIへの指示文（質問・要約方針など）:",
                  font=("", 10, "bold")).pack(anchor="w")
        instr_text = scrolledtext.ScrolledText(frm, height=5, wrap="word")
        instr_text.pack(fill="x", pady=(2, 8))
        instr_text.insert("1.0",
            "以下は社内ファイルから「{クエリ}」で検索した結果です。\n"
            "1. これらのファイルから関連情報を抽出して要約してください\n"
            "2. 最も重要なファイル上位3つを理由付きで挙げてください\n"
            "3. 矛盾や疑問点があれば指摘してください")

        # 設定
        opts = ttk.LabelFrame(frm, text="出力オプション")
        opts.pack(fill="x", pady=8)

        row1 = ttk.Frame(opts)
        row1.pack(fill="x", padx=8, pady=4)
        ttk.Label(row1, text="検索モード:").pack(side="left")
        mode_var = tk.StringVar(value=self.search_mode.get())
        ttk.Radiobutton(row1, text="AND", value="AND",
                        variable=mode_var).pack(side="left", padx=4)
        ttk.Radiobutton(row1, text="OR", value="OR",
                        variable=mode_var).pack(side="left", padx=4)

        row2 = ttk.Frame(opts)
        row2.pack(fill="x", padx=8, pady=4)
        ttk.Label(row2, text="含めるファイル数:").pack(side="left")
        n_var = tk.IntVar(value=30)
        ttk.Spinbox(row2, from_=5, to=200, textvariable=n_var,
                    width=8).pack(side="left", padx=(4, 16))
        ttk.Label(row2, text="1ファイルあたり最大文字数:").pack(side="left")
        max_chars_var = tk.IntVar(value=3000)
        ttk.Spinbox(row2, from_=200, to=20000, increment=500,
                    textvariable=max_chars_var, width=8).pack(side="left", padx=4)

        # 出力先
        out_frm = ttk.LabelFrame(frm, text="出力先")
        out_frm.pack(fill="x", pady=8)
        out_var = tk.StringVar(value="file")
        ttk.Radiobutton(out_frm, text="ファイルに保存（.md）",
                        value="file", variable=out_var).pack(anchor="w", padx=8, pady=2)
        ttk.Radiobutton(out_frm, text="クリップボードにコピー",
                        value="clip", variable=out_var).pack(anchor="w", padx=8, pady=2)

        # 推定サイズ表示
        info_var = tk.StringVar(value="")
        ttk.Label(frm, textvariable=info_var, foreground="#666").pack(anchor="w")

        def update_estimate(*_):
            est = n_var.get() * max_chars_var.get()
            tokens = est // 3  # ざっくり1トークン≒3文字
            info_var.set(
                f"推定サイズ: 最大 {est/1000:.0f}K 文字 ≒ {tokens/1000:.0f}K トークン  "
                f"(Copilot/GPT-4: 〜128K, Claude: 〜200K)"
            )
        n_var.trace_add("write", update_estimate)
        max_chars_var.trace_add("write", update_estimate)
        update_estimate()

        # 実行ボタン
        btns = ttk.Frame(frm)
        btns.pack(fill="x", pady=(12, 0))

        def do_export():
            q = query_var.get().strip()
            if not q:
                messagebox.showwarning(APP_NAME, "クエリを入力してください。", parent=win)
                return
            instruction = instr_text.get("1.0", "end").strip().replace("{クエリ}", q)
            n_limit = n_var.get()
            max_chars = max_chars_var.get()
            mode = mode_var.get()

            # 検索
            self.status_var.set("AI出力用に検索中…")
            self.update_idletasks()
            rows = self.indexer.search(q, mode=mode, limit=n_limit)
            if not rows:
                messagebox.showinfo(APP_NAME,
                                    f"検索結果が0件でした: {q}", parent=win)
                return
            # 本文一括取得
            paths = [r[0] for r in rows]
            contents_map = self.indexer.get_contents(paths)

            # Markdown生成
            md = _build_ai_markdown(q, instruction, rows, contents_map, max_chars)

            if out_var.get() == "file":
                f = filedialog.asksaveasfilename(
                    parent=win,
                    title="保存先",
                    defaultextension=".md",
                    filetypes=[("Markdown", "*.md"), ("Text", "*.txt"), ("All", "*.*")],
                    initialfile=f"ai_context_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md",
                )
                if not f:
                    return
                try:
                    Path(f).write_text(md, encoding="utf-8")
                except Exception as e:
                    messagebox.showerror(APP_NAME, f"保存失敗:\n{e}", parent=win)
                    return
                win.destroy()
                if messagebox.askyesno(
                    APP_NAME,
                    f"保存しました（{len(md):,}文字）\n\n{f}\n\nファイルを開きますか？",
                ):
                    try:
                        if sys.platform == "win32":
                            os.startfile(f)
                        elif sys.platform == "darwin":
                            subprocess.Popen(["open", f])
                        else:
                            subprocess.Popen(["xdg-open", f])
                    except Exception:
                        pass
            else:
                self.clipboard_clear()
                self.clipboard_append(md)
                self.update()  # クリップボードを確実に反映
                win.destroy()
                messagebox.showinfo(
                    APP_NAME,
                    f"クリップボードにコピーしました\n\n"
                    f"  サイズ: {len(md):,} 文字\n"
                    f"  ファイル数: {len(rows)} 件\n\n"
                    f"AIチャット画面に貼り付けて使ってください。",
                )

        ttk.Button(btns, text="出力", command=do_export).pack(side="right")
        ttk.Button(btns, text="キャンセル", command=win.destroy).pack(side="right", padx=(0, 8))

    def open_db_settings(self):
        """DBフォルダの確認・変更・デフォルト復元・旧DBからの移行のダイアログ。"""
        if self.worker and self.worker.is_alive():
            messagebox.showwarning(APP_NAME,
                "スキャン中はDB設定を変更できません。停止してから再試行してください。")
            return

        win = tk.Toplevel(self)
        win.title("DB設定")
        win.geometry("640x300")
        win.transient(self)
        win.grab_set()

        frm = ttk.Frame(win)
        frm.pack(fill="both", expand=True, padx=16, pady=12)

        # 現状表示
        ttk.Label(frm, text="現在のシャードDBフォルダ:",
                  font=("", 10, "bold")).pack(anchor="w")
        path_var = tk.StringVar(value=str(self.indexer.shard_dir))
        ttk.Entry(frm, textvariable=path_var, state="readonly",
                  font=("Consolas", 9)).pack(fill="x", pady=(2, 8))

        # シャード数とサイズの情報
        n, _ = self.indexer.stats()
        disk = 0
        shard_count = 0
        for i in range(self.indexer.num_shards):
            for suffix in ("", "-wal", "-shm"):
                p = Path(str(self.indexer.shard_path(i)) + suffix)
                if p.exists():
                    try:
                        disk += p.stat().st_size
                    except OSError:
                        pass
            if self.indexer.shard_path(i).exists():
                shard_count += 1
        info = (f"インデックス済 {n:,} ファイル / "
                f"シャードファイル {shard_count}/{self.indexer.num_shards} 個 / "
                f"合計 {disk/1e9:.2f} GB")
        ttk.Label(frm, text=info, foreground="#555").pack(anchor="w", pady=(0, 12))

        # 操作ボタン
        ttk.Label(frm, text="操作:", font=("", 10, "bold")).pack(anchor="w")

        def do_change():
            new = filedialog.askdirectory(
                title="シャードDBフォルダを選択（移動先または新規作成先）",
                initialdir=str(self.indexer.shard_dir),
                parent=win,
            )
            if not new:
                return
            new_path = Path(new)
            # 既存シャードの有無を判定
            existing_shards = sum(
                1 for i in range(NUM_SHARDS)
                if (new_path / f"index_{i:02x}.db").exists()
            )
            if existing_shards > 0:
                msg = (f"このフォルダには既に {existing_shards} 個のシャードDBがあります。\n\n"
                       f"  パス: {new_path}\n\n"
                       f"このDBを使用しますか？\n"
                       f"（はい：このDBを使用 / いいえ：キャンセル）")
            else:
                msg = (f"このフォルダにはシャードDBがありません。\n"
                       f"新規シャードを作成しますか？\n\n"
                       f"  パス: {new_path}")
            if not messagebox.askyesno(APP_NAME, msg, parent=win):
                return
            self._switch_db_dir(new_path)
            win.destroy()

        def do_reset():
            if str(self.indexer.shard_dir) == str(SHARD_DIR):
                messagebox.showinfo(APP_NAME, "既にデフォルト場所を使用中です。", parent=win)
                return
            if not messagebox.askyesno(APP_NAME,
                f"DBフォルダをデフォルトに戻しますか？\n\n  {SHARD_DIR}",
                parent=win):
                return
            self._switch_db_dir(SHARD_DIR)
            win.destroy()

        def do_migrate():
            f = filedialog.askopenfilename(
                title="移行元の旧 index.db を選択",
                filetypes=[("SQLite DB", "*.db"), ("All", "*.*")],
                parent=win,
            )
            if not f:
                return
            legacy = Path(f)
            try:
                size = legacy.stat().st_size
            except OSError:
                messagebox.showerror(APP_NAME, "ファイルが読めません。", parent=win)
                return
            if not messagebox.askyesno(APP_NAME,
                f"以下のDBから現在のシャードへ移行します:\n\n"
                f"  {legacy}\n  サイズ: {size/1e9:.2f} GB\n\n"
                f"既存のシャードデータには上書き追加されます。続行しますか？",
                parent=win):
                return
            win.destroy()
            self._run_migration_from(legacy)

        btn_row = ttk.Frame(frm)
        btn_row.pack(fill="x", pady=(8, 0))
        ttk.Button(btn_row, text="フォルダ変更…",
                   command=do_change).pack(side="left", padx=(0, 4))
        ttk.Button(btn_row, text="デフォルトに戻す",
                   command=do_reset).pack(side="left", padx=4)
        ttk.Button(btn_row, text="旧 index.db から移行…",
                   command=do_migrate).pack(side="left", padx=4)

        ttk.Button(frm, text="閉じる", command=win.destroy).pack(pady=(16, 0))

    def _switch_db_dir(self, new_dir: Path):
        """DBフォルダを切り替えて Indexer を再初期化、configに保存。"""
        new_dir = Path(new_dir).resolve()
        self.cfg["db_dir"] = str(new_dir)
        save_config(self.cfg)
        self.indexer = Indexer(new_dir)
        self._refresh_stats()
        self._log(f"DBフォルダを変更: {new_dir}")

    def _run_migration_from(self, legacy: Path):
        """指定された旧DBから現在のシャードへ移行（モーダル進捗ダイアログ付き）。"""
        dlg = tk.Toplevel(self)
        dlg.title("マイグレーション中")
        dlg.geometry("480x180")
        dlg.transient(self)
        dlg.grab_set()
        dlg.protocol("WM_DELETE_WINDOW", lambda: None)
        ttk.Label(dlg, text=f"以下のDBから現シャードへ移行中:\n{legacy}",
                  justify="center").pack(pady=10, padx=20)
        pb = ttk.Progressbar(dlg, length=420, mode="indeterminate")
        pb.pack(pady=5, padx=20)
        pb.start(10)
        lbl = ttk.Label(dlg, text="旧DBを開いています…")
        lbl.pack()

        mq: queue.Queue = queue.Queue()
        def worker():
            try:
                def cb(done, total):
                    mq.put(("progress", done, total))
                count = migrate_from_legacy(legacy, self.indexer, cb)
                mq.put(("done", count, None))
            except Exception as e:
                mq.put(("error", str(e), None))
        threading.Thread(target=worker, daemon=True).start()

        started = [False]
        def poll():
            try:
                while True:
                    kind, a, b = mq.get_nowait()
                    if kind == "progress":
                        if b > 0 and not started[0]:
                            pb.stop()
                            pb.configure(mode="determinate", maximum=b)
                            started[0] = True
                        if b > 0:
                            pb["value"] = a
                            lbl["text"] = f"{a:,} / {b:,} 件 ({a*100//max(b,1)}%)"
                        else:
                            lbl["text"] = "メタ情報を読み込み中…"
                    elif kind == "done":
                        pb.stop()
                        dlg.destroy()
                        self._refresh_stats()
                        self._log(f"移行完了: {a:,} 件 ({legacy})")
                        messagebox.showinfo(APP_NAME, f"移行完了: {a:,} 件")
                        return
                    elif kind == "error":
                        pb.stop()
                        dlg.destroy()
                        messagebox.showerror(APP_NAME, f"移行失敗:\n{a}")
                        return
            except queue.Empty:
                pass
            self.after(100, poll)
        poll()

    def _refresh_stats(self):
        n, total = self.indexer.stats()
        # シャードのディスク使用量を合計
        disk = 0
        for i in range(self.indexer.num_shards):
            for suffix in ("", "-wal", "-shm"):
                p = Path(str(self.indexer.shard_path(i)) + suffix)
                if p.exists():
                    try:
                        disk += p.stat().st_size
                    except OSError:
                        pass
        # 長いパスは短縮
        db_path = str(self.indexer.shard_dir)
        if len(db_path) > 60:
            db_path = "…" + db_path[-57:]
        self.stats_var.set(
            f"インデックス済 {n:,} / {self.indexer.num_shards}シャード / "
            f"{disk/1e9:.2f}GB  [DB: {db_path}]"
        )

    # ---- 旧DBからのマイグレーション ----
    def _maybe_migrate_legacy(self):
        if not LEGACY_DB.exists():
            return
        try:
            legacy_size = LEGACY_DB.stat().st_size
        except OSError:
            return
        # 既にシャードにデータがある場合はスキップ
        n, _ = self.indexer.stats()
        if n > 0:
            self._log(f"※ 旧DB {LEGACY_DB.name} がありますが、"
                      f"シャードに既にデータがあるため移行をスキップしました。")
            return
        ok = messagebox.askyesno(
            APP_NAME,
            f"旧形式の単一DB を検出しました\n\n"
            f"  パス: {LEGACY_DB}\n"
            f"  サイズ: {legacy_size/1e9:.2f} GB\n\n"
            f"{self.indexer.num_shards} 個のシャードDBに分割マイグレーションしますか？\n"
            f"（10GBの場合 数分〜十数分かかる可能性があります）\n\n"
            f"「いいえ」を選ぶと旧DBは無視され、シャードを新規作成して使います。"
        )
        if not ok:
            return
        self._run_migration()

    def _run_migration(self):
        """モーダルダイアログでマイグレーション進行表示。"""
        dlg = tk.Toplevel(self)
        dlg.title("マイグレーション中")
        dlg.geometry("480x180")
        dlg.transient(self)
        dlg.grab_set()
        dlg.protocol("WM_DELETE_WINDOW", lambda: None)

        ttk.Label(
            dlg, text="旧DBから分割シャードDBへデータ移行中…\n（中断不可・ウィンドウを閉じないでください）",
            justify="center"
        ).pack(pady=10, padx=20)
        pb = ttk.Progressbar(dlg, length=420, mode="indeterminate")
        pb.pack(pady=5, padx=20)
        pb.start(10)  # 進捗カウントが来るまでは不定形アニメ
        lbl = ttk.Label(dlg, text="旧DBを開いています…")
        lbl.pack()
        sub = ttk.Label(dlg, text="（10GBクラスでもメタ情報の読み込みに数十秒〜数分かかります）",
                        foreground="#666", font=("", 9))
        sub.pack()

        mq: queue.Queue = queue.Queue()

        def worker():
            try:
                def cb(done, total):
                    mq.put(("progress", done, total))
                count = migrate_from_legacy(LEGACY_DB, self.indexer, cb)
                mq.put(("done", count, None))
            except Exception as e:
                mq.put(("error", str(e), None))

        threading.Thread(target=worker, daemon=True).start()

        started = [False]
        def poll():
            try:
                while True:
                    kind, a, b = mq.get_nowait()
                    if kind == "progress":
                        if b > 0 and not started[0]:
                            # 総数確定 → 進捗バーを determinate に切り替え
                            pb.stop()
                            pb.configure(mode="determinate", maximum=b)
                            started[0] = True
                        if b > 0:
                            pb["value"] = a
                            lbl["text"] = f"{a:,} / {b:,} 件 ({a*100//max(b,1)}%)"
                        else:
                            lbl["text"] = "メタ情報を読み込み中…"
                    elif kind == "done":
                        pb.stop()
                        dlg.destroy()
                        bak = LEGACY_DB.with_suffix(".db.migrated.bak")
                        try:
                            if bak.exists():
                                bak.unlink()
                            LEGACY_DB.rename(bak)
                        except Exception:
                            pass
                        self._refresh_stats()
                        self._log(f"マイグレーション完了: {a:,} 件移行 → {bak.name} にリネーム")
                        messagebox.showinfo(
                            APP_NAME,
                            f"マイグレーション完了\n\n"
                            f"  移行件数: {a:,} 件\n"
                            f"  旧DB: {bak.name} としてバックアップ\n\n"
                            f"不要なら手動で削除してOK。"
                        )
                        return
                    elif kind == "error":
                        pb.stop()
                        dlg.destroy()
                        messagebox.showerror(APP_NAME, f"マイグレーション失敗:\n{a}")
                        return
            except queue.Empty:
                pass
            self.after(100, poll)
        poll()

    # ---- 検索 ----
    def do_search(self):
        q = self.query_var.get().strip()
        self.results.delete(*self.results.get_children())
        self._results_data.clear()
        self.preview.delete("1.0", "end")
        if not q:
            return
        mode = self.search_mode.get()
        t0 = time.time()
        rows = self.indexer.search(q, mode=mode)
        for path, name, snip in rows:
            iid = self.results.insert("", "end", text=name, values=(path,))
            self._results_data.append((iid, path, snip))
        self.status_var.set(
            f"{len(rows)} 件ヒット ({(time.time()-t0)*1000:.0f} ms) [モード:{mode}]"
        )

    def on_select(self, _):
        sel = self.results.selection()
        if not sel:
            return
        iid = sel[0]
        for rec in self._results_data:
            if rec[0] == iid:
                _, path, snip = rec
                self.preview.delete("1.0", "end")
                self.preview.insert("end", path + "\n", "path")
                self.preview.insert("end", "─" * 60 + "\n")
                # スニペット内の 〘...〙 をハイライト
                idx = "3.0"
                self.preview.insert("end", snip + "\n\n")
                self._highlight_markers()

                # 余裕があればファイルの中身も読み込んで該当箇所を多めに表示
                self._show_more_hits(path)
                break

    def _highlight_markers(self):
        # 〘...〙 をハイライト
        text = self.preview.get("1.0", "end")
        for m in re.finditer(r"〘([^〙]*)〙", text):
            start = f"1.0+{m.start()}c"
            end = f"1.0+{m.end()}c"
            self.preview.tag_add("hit", start, end)

    def _show_more_hits(self, path: str):
        """ファイルを開きなおして、検索語の周辺も表示する。"""
        q = self.query_var.get().strip()
        if not q:
            return
        terms = [t.strip('"') for t in re.findall(r'"[^"]+"|\S+', q) if not t.startswith("-")]
        if not terms:
            return
        try:
            text = extract_text(Path(path)) or ""
        except Exception:
            return
        if not text:
            return

        self.preview.insert("end", "── ファイル内のヒット箇所 ──\n\n")
        pattern = re.compile("|".join(re.escape(t) for t in terms), re.IGNORECASE)
        shown = 0
        for m in pattern.finditer(text):
            if shown >= 20:
                self.preview.insert("end", "（...以下省略）\n")
                break
            s = max(0, m.start() - 80)
            e = min(len(text), m.end() + 80)
            ctx = text[s:e].replace("\n", " ")
            insert_pos = self.preview.index("end")
            self.preview.insert("end", f"…{ctx}…\n\n")
            # ハイライト
            for mm in pattern.finditer(ctx):
                start = f"{insert_pos}+{1 + mm.start()}c"
                end = f"{insert_pos}+{1 + mm.end()}c"
                self.preview.tag_add("hit", start, end)
            shown += 1

    def on_open(self, _=None):
        sel = self.results.selection()
        if not sel:
            return
        path = self.results.item(sel[0], "values")[0]
        try:
            if sys.platform == "win32":
                os.startfile(path)
            elif sys.platform == "darwin":
                subprocess.Popen(["open", path])
            else:
                subprocess.Popen(["xdg-open", path])
        except Exception as e:
            messagebox.showerror(APP_NAME, f"開けませんでした:\n{e}")

    def open_location(self, _=None):
        """エクスプローラで該当ファイルを選択した状態で親フォルダを開く。"""
        sel = self.results.selection()
        if not sel:
            return
        path = self.results.item(sel[0], "values")[0]
        try:
            if sys.platform == "win32":
                # /select でファイルを選択状態にして開く
                # explorer はカンマ後のスペースでも動くがクオートを正しく処理させるため
                # Popen のリスト渡しでなく文字列で渡す（パスにスペースがあっても OK）
                subprocess.Popen(f'explorer /select,"{path}"')
            elif sys.platform == "darwin":
                subprocess.Popen(["open", "-R", path])
            else:
                # Linux: 親フォルダを開く（ファイル選択は環境依存）
                parent = str(Path(path).parent)
                subprocess.Popen(["xdg-open", parent])
        except Exception as e:
            messagebox.showerror(APP_NAME, f"場所を開けませんでした:\n{e}")

    def copy_path(self, _=None):
        sel = self.results.selection()
        if not sel:
            return
        path = self.results.item(sel[0], "values")[0]
        self.clipboard_clear()
        self.clipboard_append(path)
        self.status_var.set(f"パスをコピー: {path}")

    def copy_name(self, _=None):
        sel = self.results.selection()
        if not sel:
            return
        name = self.results.item(sel[0], "text")
        self.clipboard_clear()
        self.clipboard_append(name)
        self.status_var.set(f"ファイル名をコピー: {name}")

    def _show_ctx_menu(self, event):
        # 右クリックされた行を選択状態に
        iid = self.results.identify_row(event.y)
        if iid:
            self.results.selection_set(iid)
            try:
                self.ctx_menu.tk_popup(event.x_root, event.y_root)
            finally:
                self.ctx_menu.grab_release()

    # ---- 進捗キュー ----
    def _poll_queue(self):
        try:
            while True:
                kind, payload = self.q.get_nowait()
                if kind == "log":
                    self._log(payload)
                elif kind == "status":
                    self.status_var.set(payload)
                elif kind == "error":
                    self._log("エラー:\n" + payload)
                elif kind == "done":
                    self.btn_index.config(state="normal")
                    self.btn_full.config(state="normal")
                    self.btn_stop.config(state="disabled")
                    self.progress.stop()
                    self._refresh_stats()
                    self.status_var.set("インデックス更新完了")
        except queue.Empty:
            pass

        # スキャン中は1秒ごとに「インデックス済」カウントをDBから更新
        if self.worker and self.worker.is_alive():
            now = time.time()
            if now - self._last_stats_update > 1.0:
                self._refresh_stats()
                self._last_stats_update = now

        self.after(50, self._poll_queue)

    def _log(self, msg: str):
        ts = datetime.now().strftime("%H:%M:%S")
        self.log.insert("end", f"[{ts}] {msg}\n")
        self.log.see("end")


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------

def main():
    app = App()
    app.mainloop()


if __name__ == "__main__":
    main()
